"""PPTX presentation content extractor using python-pptx."""

from pathlib import Path
from typing import Generator, Union

import pptx

from src.ingestion.detectors.base import DetectionResult, FileType
from src.ingestion.extractors.base import BaseExtractor
from src.ingestion.extractors.models import ContentType, ExtractedElement, ExtractedStream, ExtractionStatus


class PPTXExtractor(BaseExtractor):
    """Extracts slide text, titles, and tables from PPTX files using python-pptx."""

    def __init__(self):
        super().__init__(target_file_type=FileType.PPTX)

    def _extract_implementation(self, path: Path, detection_result: DetectionResult) -> ExtractedStream:
        try:
            prs = pptx.Presentation(path)
            total_slides = len(prs.slides)
            metadata = {
                "slide_count": total_slides,
            }

            generator = self._slide_generator(prs)
            return ExtractedStream(
                file_path=str(path),
                file_type=FileType.PPTX,
                element_generator=generator,
                metadata=metadata,
            )

        except Exception as e:
            return ExtractedStream(
                file_path=str(path),
                file_type=FileType.PPTX,
                element_generator=iter([]),
                status=ExtractionStatus.FAILED,
                error_message=f"Failed to open PPTX presentation: {str(e)}",
            )

    def _slide_generator(self, prs: pptx.Presentation) -> Generator[ExtractedElement, None, None]:
        element_idx = 0
        total_slides = len(prs.slides)

        for s_idx, slide in enumerate(prs.slides):
            slide_number = s_idx + 1
            slide_title = None

            if slide.shapes.title and slide.shapes.title.text:
                slide_title = slide.shapes.title.text.strip()

            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
                elif shape.has_table:
                    for row in shape.table.rows:
                        row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_cells:
                            slide_texts.append(" | ".join(row_cells))

            combined_text = "\n".join(slide_texts).strip()
            if not combined_text:
                continue

            yield ExtractedElement(
                element_index=element_idx,
                content_type=ContentType.SLIDE,
                text_content=combined_text,
                structural_context={
                    "slide_number": slide_number,
                    "total_slides": total_slides,
                    "slide_title": slide_title,
                },
            )
            element_idx += 1
